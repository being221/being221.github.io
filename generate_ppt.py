from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ============================================================
# 主题配色
# ============================================================
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
ACCENT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF4, 0xF4)
TEXT_DARK = RGBColor(0x2C, 0x3E, 0x50)
TEXT_BODY = RGBColor(0x34, 0x49, 0x5E)

prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    """填充纯色背景"""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, color, opacity=None):
    """添加矩形色块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=TEXT_BODY, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    """添加文本框"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_title_bar(slide, title_text):
    """添加统一的顶部标题栏"""
    add_rect(slide, 0, 0, W, Inches(1.3), DARK_BLUE)
    # 底部细线
    add_rect(slide, 0, Inches(1.3), W, Inches(0.04), ACCENT_BLUE)
    add_text_box(slide, Inches(1), Inches(0.25), Inches(11), Inches(0.9),
                 title_text, font_size=36, color=WHITE, bold=True)


def add_bullet_list(slide, left, top, width, height, items, font_size=20):
    """添加项目符号列表"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = TEXT_BODY
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(12)
        p.level = 0
    return txBox

# ============================================================
# 第1页：封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白
add_bg(slide, DARK_BLUE)
# 装饰线
add_rect(slide, Inches(1.5), Inches(2.8), Inches(1.2), Inches(0.05), ACCENT_BLUE)

add_text_box(slide, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
             "毕业实习汇报", font_size=52, color=WHITE, bold=True)
add_text_box(slide, Inches(1.5), Inches(3.1), Inches(10), Inches(0.8),
             "钢筋混凝土框架结构设计", font_size=28, color=ACCENT_BLUE)
add_text_box(slide, Inches(1.5), Inches(4.3), Inches(8), Inches(0.5),
             "土木工程专业  |  22141010104  |  邓杰鹏", font_size=20, color=RGBColor(0xBD, 0xC3, 0xCA))
add_text_box(slide, Inches(1.5), Inches(4.9), Inches(8), Inches(0.5),
             "实习单位：XX建筑设计研究院  |  指导老师：XXX", font_size=18, color=RGBColor(0x95, 0xA5, 0xA6))

# ============================================================
# 第2页：目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "目录 CONTENTS")
items = [
    "01   实习概况",
    "02   实习目的与意义",
    "03   实习内容与任务",
    "04   结构设计流程",
    "05   主要成果展示",
    "06   实习收获与体会",
    "07   总结与致谢",
]
add_bullet_list(slide, Inches(3), Inches(2.0), Inches(8), Inches(5), items, font_size=26)

# ============================================================
# 第3页：实习概况
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "01  实习概况")
items = [
    "实习时间：2025年12月 — 2026年4月（共16周）",
    "实习单位：XX建筑设计研究院有限公司",
    "实习岗位：结构设计实习生",
    "实习项目：某多层钢筋混凝土框架结构办公楼设计",
    "项目概况：建筑总高度17.7m，共5层，建筑面积约4200m²",
    "结构体系：钢筋混凝土框架结构，抗震设防烈度7度（0.15g）",
    "主要工作：结构方案选型、荷载计算、内力分析、截面设计、施工图绘制",
]
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), items, font_size=22)

# ============================================================
# 第4页：实习目的与意义
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "02  实习目的与意义")

# 左侧
add_rect(slide, Inches(0.8), Inches(1.8), Inches(5.5), Inches(0.6), ACCENT_BLUE)
add_text_box(slide, Inches(1.0), Inches(1.85), Inches(5), Inches(0.5),
             "实习目的", font_size=24, color=WHITE, bold=True)
items_left = [
    "将理论知识与工程实践相结合",
    "掌握结构设计全流程操作方法",
    "熟悉国家现行结构设计规范",
    "培养独立解决工程问题的能力",
    "了解设计院工作流程与协作方式",
]
add_bullet_list(slide, Inches(1.0), Inches(2.6), Inches(5.2), Inches(4), items_left, font_size=20)

# 右侧
add_rect(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(0.6), ACCENT_BLUE)
add_text_box(slide, Inches(7.2), Inches(1.85), Inches(5), Inches(0.5),
             "实习意义", font_size=24, color=WHITE, bold=True)
items_right = [
    "为毕业设计积累完整工程经验",
    "提升CAD、计算软件实操技能",
    "建立结构工程师的职业认知",
    "培养严谨细致的工作态度",
    "为就业面试提供项目经历支撑",
]
add_bullet_list(slide, Inches(7.2), Inches(2.6), Inches(5.2), Inches(4), items_right, font_size=20)

# ============================================================
# 第5页：实习内容与任务
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "03  实习内容与任务")
items = [
    "结构方案选型 — 根据建筑功能与跨度确定框架结构体系，初选梁柱截面尺寸",
    "荷载统计与计算 — 完成恒载、活载、风荷载、地震作用的全面统计",
    "内力分析与组合 — 采用D值法、弯矩分配法进行水平与竖向荷载内力计算",
    "截面设计与配筋 — 按《混凝土结构设计规范》完成梁柱板的配筋设计",
    "施工图绘制 — 使用CAD绘制全套结构施工图（平面、立面、节点详图）",
    "计算书编制 — 整理完整的设计计算书，包含全部计算过程与结果校核",
]
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), items, font_size=21)

# ============================================================
# 第6页：结构设计流程
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "04  结构设计流程")

steps = [
    ("1", "结构方案\n确定", "选定框架结构\n初选构件截面"),
    ("2", "荷载\n统计", "恒/活/风/震\n全面统计"),
    ("3", "内力\n计算", "D值法+弯矩\n分配法"),
    ("4", "截面\n设计", "梁柱板配筋\n规范校核"),
    ("5", "图纸\n绘制", "全套结构\n施工图"),
]
for i, (num, title, desc) in enumerate(steps):
    x = Inches(1.2 + i * 2.35)
    y_top = Inches(2.2)
    # 圆形编号
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.6), y_top, Inches(0.8), Inches(0.8))
    circle.fill.solid()
    circle.fill.fore_color.rgb = ACCENT_BLUE
    circle.line.fill.background()
    tf = circle.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = num
    p.font.size = Pt(24)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.font.name = 'Microsoft YaHei'
    p.alignment = PP_ALIGN.CENTER
    # 箭头
    if i < len(steps) - 1:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(1.6), y_top + Inches(0.25), Inches(0.55), Inches(0.3))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = RGBColor(0xBD, 0xC3, 0xCA)
        arrow.line.fill.background()
    # 标题
    add_text_box(slide, x + Inches(0.1), y_top + Inches(1.0), Inches(1.8), Inches(0.6),
                 title, font_size=20, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    # 描述
    add_text_box(slide, x + Inches(0.1), y_top + Inches(1.6), Inches(1.8), Inches(0.8),
                 desc, font_size=14, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

# 底部说明
add_text_box(slide, Inches(1), Inches(5.2), Inches(11), Inches(1.5),
             "结构设计遵循\"方案确定 → 荷载统计 → 内力分析 → 截面设计 → 构造措施\"的标准流程，"
             "各环节相互关联、逐步深化，前一环节的输出即为后一环节的输入依据。",
             font_size=16, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第7页：主要成果展示
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "05  主要成果展示")

# 四个成果卡片
cards = [
    ("计算书", "完成完整的结构设计\n计算书，含荷载统计、\n内力分析、配筋计算\n共6章约80页"),
    ("结构施工图", "绘制全套结构施工图\n含平面图、立面图、\n节点详图共11张"),
    ("模型数据", "建立PKPM结构分析\n模型，完成整体指标\n校核（位移、轴压比）"),
    ("设计总结", "撰写实习报告与\n设计总结，归纳关键\n技术要点与经验"),
]
for i, (title, desc) in enumerate(cards):
    x = Inches(0.8 + i * 3.1)
    y = Inches(2.0)
    card = add_rect(slide, x, y, Inches(2.8), Inches(4.2), LIGHT_GRAY)
    add_rect(slide, x, y, Inches(2.8), Inches(0.06), ACCENT_BLUE)
    add_text_box(slide, x + Inches(0.3), y + Inches(0.3), Inches(2.2), Inches(0.5),
                 title, font_size=24, color=DARK_BLUE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.3), y + Inches(1.1), Inches(2.2), Inches(2.8),
                 desc, font_size=16, color=TEXT_BODY, alignment=PP_ALIGN.CENTER)

# ============================================================
# 第8页：实习收获与体会
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_title_bar(slide, "06  实习收获与体会")

items = [
    "理论到实践的跨越 — 真正体会了\"按规范设计\"与\"按教材计算\"的差异，理解了安全度、经济性、可施工性三者之间的权衡",
    "规范应用能力提升 — 系统学习了《建筑结构荷载规范》《混凝土结构设计规范》《建筑抗震设计规范》在实际工程中的综合运用",
    "软件技能进阶 — CAD绘图效率显著提高，掌握了PKPM建模与结果分析，学会了用Python辅助批量计算与数据处理",
    "工程思维养成 — 认识到结构设计不是孤立的技术问题，需要与建筑、设备专业协调，兼顾业主需求与施工可行性",
    "职业认知深化 — 通过参与真实项目，对结构工程师的职业路径、能力要求、行业现状有了更清晰的认识",
]
add_bullet_list(slide, Inches(1), Inches(1.8), Inches(11), Inches(5), items, font_size=20)

# ============================================================
# 第9页：总结与致谢
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_rect(slide, Inches(2), Inches(2.0), Inches(1.2), Inches(0.05), ACCENT_BLUE)

add_text_box(slide, Inches(2), Inches(1.0), Inches(9), Inches(0.8),
             "总结与致谢", font_size=44, color=WHITE, bold=True)

summary = (
    "本次毕业实习以\"钢筋混凝土框架结构设计\"为核心任务，历经16周的工程实践，"
    "完整参与了从结构方案确定到施工图绘制的全流程。通过实操训练，"
    "不仅巩固了专业理论知识，更建立了工程实践的思维方式，为即将开始的职业生涯奠定了坚实基础。\n\n"
    "衷心感谢实习单位提供的宝贵实践平台，感谢指导老师在专业技术上的悉心指导，"
    "也感谢实习期间各位同事的帮助与支持。"
)
add_text_box(slide, Inches(2), Inches(2.3), Inches(9), Inches(3.5),
             summary, font_size=22, color=RGBColor(0xBD, 0xC3, 0xCA))

add_text_box(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5),
             "谢谢！", font_size=40, color=ACCENT_BLUE, bold=True)

# ============================================================
# 保存
# ============================================================
output = r"C:\Users\邓杰鹏\Desktop\毕设\毕设\毕业实习汇报.pptx"
prs.save(output)
print(f"已生成: {output}")
print(f"共 {len(prs.slides)} 页")
