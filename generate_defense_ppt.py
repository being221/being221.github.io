# -*- coding: utf-8 -*-
"""
生成毕设答辩PPT - 学术严谨风格 - 数据全部来自计算书原件
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ── 颜色方案 ──
DARK_BLUE = RGBColor(0x1B, 0x3A, 0x5C)
LIGHT_BLUE = RGBColor(0x2E, 0x86, 0xC1)
BLACK = RGBColor(0x00, 0x00, 0x00)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRAY_BG = RGBColor(0xF2, 0xF4, 0xF4)
DARK_GRAY = RGBColor(0x56, 0x64, 0x6E)
TABLE_HEADER_BG = RGBColor(0x1B, 0x3A, 0x5C)
TABLE_ALT_BG = RGBColor(0xEB, 0xF5, 0xFB)

prs = Presentation()
prs.slide_width = Cm(33.867)  # 16:9
prs.slide_height = Cm(19.05)

SLD_W = 33.867
SLD_H = 19.05


def add_blank_slide():
    layout = prs.slide_layouts[6]
    return prs.slides.add_slide(layout)


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                font_color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="微软雅黑", line_spacing=1.2):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    p.space_after = Pt(line_spacing * font_size - font_size)
    for run in p.runs:
        rPr = run._r.get_or_add_rPr()
        rPr.set(qn('a:eaTypeface'), font_name)
    return tf


def add_line(slide, x1, y1, x2, y2, color=DARK_BLUE, width=1.5):
    connector = slide.shapes.add_connector(
        1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    connector.line.color.rgb = color
    connector.line.width = Pt(width)
    return connector


def add_title_bar(slide, title_text, subtitle_text=""):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(SLD_W), Cm(2.8))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE
    shape.line.fill.background()
    add_textbox(slide, 1.5, 0.4, 30, 1.5, title_text,
                font_size=26, font_color=WHITE, bold=True)
    add_line(slide, 1.5, 2.9, SLD_W - 1.5, 2.9, LIGHT_BLUE, 1)
    if subtitle_text:
        add_textbox(slide, 1.5, 1.7, 30, 0.8, subtitle_text,
                    font_size=12, font_color=RGBColor(0xBB, 0xCC, 0xDD))


def add_page_number(slide, num, total=18):
    add_textbox(slide, SLD_W - 3, SLD_H - 1.2, 2.5, 0.8,
                f"{num}/{total}", font_size=10, font_color=DARK_GRAY,
                alignment=PP_ALIGN.RIGHT)


def make_table(slide, left, top, col_widths, headers, rows, font_size=10):
    n_rows = len(rows) + 1
    n_cols = len(headers)
    tbl_height = 0.65 * n_rows
    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Cm(left), Cm(top),
        Cm(sum(col_widths)), Cm(tbl_height))
    table = table_shape.table
    for i, w in enumerate(col_widths):
        table.columns[i].width = Cm(w)
    # 表头
    for j, hdr in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "微软雅黑"
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                rPr = run._r.get_or_add_rPr()
                rPr.set(qn('a:eaTypeface'), "微软雅黑")
    # 数据行
    for i, row_data in enumerate(rows):
        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.name = "微软雅黑"
                p.alignment = PP_ALIGN.CENTER
                for run in p.runs:
                    rPr = run._r.get_or_add_rPr()
                    rPr.set(qn('a:eaTypeface'), "微软雅黑")
    return table_shape


def add_bullet_list(slide, left, top, width, items, font_size=13,
                    color=BLACK, spacing=1.5):
    txBox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(12))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "微软雅黑"
        p.space_after = Pt(spacing * font_size - font_size)
        p.level = 0
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        buChar = pPr.makeelement(qn('a:buChar'), {'char': '•'})
        pPr.append(buChar)
        for run in p.runs:
            rPr = run._r.get_or_add_rPr()
            rPr.set(qn('a:eaTypeface'), "微软雅黑")
    return tf


# ═══════════════════════════════════════════════════════════
# 幻灯片1：封面
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(SLD_W), Cm(SLD_H))
shape.fill.solid()
shape.fill.fore_color.rgb = RGBColor(0xF8, 0xFA, 0xFC)
shape.line.fill.background()

shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(SLD_W), Cm(7.5))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

add_textbox(slide, 2, 1.5, 30, 1.2,
            "四川轻化工大学  土木工程学院",
            font_size=18, font_color=RGBColor(0xBB, 0xCC, 0xDD))

add_textbox(slide, 2, 3.2, 30, 2.0,
            "多层钢筋混凝土框架结构设计",
            font_size=40, font_color=WHITE, bold=True)
add_textbox(slide, 2, 5.2, 30, 1.0,
            "—— 海心小区7栋住宅楼结构计算",
            font_size=20, font_color=RGBColor(0xDD, 0xE8, 0xF0))

info_items = [
    "答辩人：邓杰鹏",
    "学  号：22141010104",
    "专  业：土木工程 2021级2班",
    "指导教师：赵蕴林",
]
for i, item in enumerate(info_items):
    add_textbox(slide, 2, 9.0 + i * 1.2, 15, 1.0, item,
                font_size=16, font_color=DARK_BLUE, bold=(i == 3))

add_textbox(slide, 2, 16.0, 10, 0.8, "2025年6月",
            font_size=14, font_color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════
# 幻灯片2：目录
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "目  录")
add_page_number(slide, 2)

toc_items = [
    ("1", "工程概况与设计依据", "项目背景、材料选用、规范依据"),
    ("2", "结构布置与截面尺寸", "计算简图、梁柱截面初选"),
    ("3", "侧向刚度计算（D值法）", "框架整体抗侧刚度"),
    ("4", "竖向荷载统计", "恒载、活载导算"),
    ("5", "水平地震作用下内力计算", "底部剪力法、反弯点法"),
    ("6", "风荷载作用下内力计算", "水平风荷载及层间位移"),
    ("7", "竖向荷载下内力计算", "弯矩二次分配法"),
    ("8", "内力组合", "无震组合 + 有震组合"),
    ("9", "截面设计", "梁柱正截面+斜截面配筋"),
    ("10", "楼板与楼梯设计", "双向板弹性理论、板式楼梯"),
    ("11", "基础设计", "独立基础 + 联合基础"),
    ("12", "结论与致谢", "设计总结"),
]
y_start = 3.8
for i, (num, title, desc) in enumerate(toc_items):
    y = y_start + i * 1.2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, Cm(2.5), Cm(y + 0.1), Cm(0.7), Cm(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_BLUE if i < 11 else LIGHT_BLUE
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].font.name = "微软雅黑"

    add_textbox(slide, 3.6, y, 15, 0.6, title,
                font_size=14, font_color=DARK_BLUE, bold=True)
    add_textbox(slide, 18, y + 0.05, 14, 0.6, desc,
                font_size=11, font_color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════
# 幻灯片3：工程概况
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "工程概况", "第1章  项目基本信息")
add_page_number(slide, 3)

info_left = [
    "项目名称：海心小区7栋住宅设计",
    "建设地点：四川省内江市",
    "建筑层数：共6层，无地下室",
    "结构形式：现浇钢筋混凝土框架结构",
    "建筑总高度：18.6m",
    "基础形式：柱下独立基础",
    "总建筑面积：3345.75 m²",
]
info_right = [
    "建筑场地类别：Ⅱ类",
    "设计使用年限：50年",
    "建筑耐火等级：二级",
    "结构安全等级：二级",
    "抗震设防烈度：7度（0.15g），第一组，Tg=0.35s",
    "框架抗震等级：三级",
    "基本风压：0.30 kN/m²   基本雪压：0.35 kN/m²",
]
add_bullet_list(slide, 1.5, 3.8, 15, info_left, font_size=13)
add_bullet_list(slide, 18, 3.8, 15, info_right, font_size=13)


# ═══════════════════════════════════════════════════════════
# 幻灯片4：结构布置与截面尺寸
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "结构布置与截面尺寸", "第2章  ④轴框架计算简图")
add_page_number(slide, 4)

add_textbox(slide, 1.5, 3.5, 30, 1.0,
            "取④轴横向框架进行手算。底层计算高度取4m（3m层高+1m基础顶面至室内地坪），其余各层取3m。"
            "各层层高均为3m，室外地坪标高-0.6m。",
            font_size=13, font_color=BLACK)

headers = ["构件", "截面尺寸 (mm)", "跨度 (mm)", "混凝土等级"]
rows = [
    ["横向边跨梁", "250×500", "5400", "C30"],
    ["横向中跨梁", "250×400", "2400", "C30"],
    ["纵向框架梁", "250×550", "6900", "C30"],
    ["次梁", "200×400", "—", "C30"],
    ["边柱 / 中柱", "500×500", "—", "C30"],
    ["楼面板 / 屋面板", "120（厚）", "—", "C30"],
]
make_table(slide, 1.5, 5.0, [5, 5.5, 4, 4], headers, rows, font_size=11)

add_textbox(slide, 1.5, 11.0, 30, 2.5,
            "柱网布置：横向三跨（AB=5400, BC=2400, CD=5400）× 纵向七跨（各6900），一字型布局，对称规整。\n"
            "柱截面由轴压比限值反算确定，N = 1.3×12×4.55×6×1000/(14.3×0.85×100)，"
            "得最小边长428.9mm，经综合考虑取500×500。",
            font_size=12, font_color=BLACK)


# ═══════════════════════════════════════════════════════════
# 幻灯片5：设计依据与材料选用
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "设计依据与材料选用", "第1章  材料选取")
add_page_number(slide, 5)

add_textbox(slide, 1.5, 3.5, 15, 0.8, "主要设计规范", font_size=16,
            font_color=DARK_BLUE, bold=True)
norms = [
    "GB50068-2018  建筑结构可靠性设计统一标准",
    "GB50009-2012  建筑结构荷载规范",
    "GB/T50011-2010  建筑抗震设计规范（2024版）",
    "GB/T50010-2010  混凝土结构设计规范（2024版）",
    "GB50007-2011  建筑地基基础设计规范",
    "GB55001-2021  工程结构通用规范",
    "GB55008-2021  混凝土结构通用规范",
]
add_bullet_list(slide, 1.5, 4.3, 15, norms, font_size=11, spacing=1.2)

add_textbox(slide, 18, 3.5, 15, 0.8, "材料参数", font_size=16,
            font_color=DARK_BLUE, bold=True)
headers = ["材料", "参数", "取值"]
rows = [
    ["混凝土 C30", "fc（轴心抗压）", "14.3 N/mm²"],
    ["混凝土 C30", "ft（轴心抗拉）", "1.43 N/mm²"],
    ["钢筋 HRB400", "fy = fyv", "360 N/mm²"],
    ["加气混凝土砌块", "容重", "7.5 kN/m³"],
    ["外墙自重", "20mm水泥+240砌块+20mm抹灰", "2.50 kN/m²"],
    ["内墙自重", "20mm抹灰+240砌块+20mm抹灰", "2.48 kN/m²"],
    ["女儿墙自重", "水泥粉刷+240墙体+20mm抹灰", "5.10 kN/m²"],
]
make_table(slide, 18, 4.5, [4, 6, 4], headers, rows, font_size=11)


# ═══════════════════════════════════════════════════════════
# 幻灯片6：D值法侧向刚度
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "侧向刚度计算（D值法）", "第2章  框架侧向刚度")
add_page_number(slide, 6)

add_textbox(slide, 1.5, 3.5, 30, 0.8,
            "采用D值法计算框架各层抗侧刚度，梁柱线刚度经计算后查表确定αc修正系数。",
            font_size=13)

headers = ["柱类型", "楼层", "K", "αc", "D (kN/m)"]
rows = [
    ["边柱", "标准层（2-6层）", "0.56", "0.22", "15097"],
    ["中柱", "标准层（2-6层）", "1.20", "0.37", "25981"],
    ["边柱", "首层", "0.74", "0.45", "13263"],
    ["中柱", "首层", "1.59", "0.58", "17070"],
]
make_table(slide, 1.5, 5.0, [4, 5, 4, 4, 4.5], headers, rows, font_size=12)

add_textbox(slide, 1.5, 9.5, 30, 4.0,
            "框架整体D值汇总（边榀2榀+中间榀5榀，共7榀）：\n"
            "• 标准层总D值 = 2×66940 + 5×82156 = 544658 kN/m\n"
            "• 首层总D值 = 2×55291 + 5×60665 = 413908 kN/m\n\n"
            "梁线刚度（中榀，考虑楼板翼缘作用）：\n"
            "• 边跨梁 ib = 2.48×10⁴ kN·m\n"
            "• 中跨梁 ib = 3.33×10⁴ kN·m\n\n"
            "柱线刚度：标准层 ic = 5.21×10⁴ kN·m，首层 ic = 3.91×10⁴ kN·m",
            font_size=12)


# ═══════════════════════════════════════════════════════════
# 幻灯片7：竖向荷载统计
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "竖向荷载统计", "第3章  恒载与活载导算")
add_page_number(slide, 7)

headers = ["荷载类型", "部位", "标准值 (kN/m²)"]
rows = [
    ["屋面恒载", "SBS防水+40找平+100珍珠岩+120板+15抹灰", "4.95"],
    ["楼面恒载", "14水磨石+30找平+120板+15抹灰", "4.20"],
    ["楼面活载", "楼面", "2.00"],
    ["走廊活载", "走廊", "2.00"],
    ["屋面活载", "屋面（不上人）", "0.50"],
    ["楼梯间活载", "楼梯间", "3.50"],
]
make_table(slide, 1.5, 4.0, [4, 9, 4], headers, rows, font_size=12)

add_textbox(slide, 1.5, 10.0, 30, 5.0,
            "导荷方式：采用双向板梯形/三角形导荷，梯形荷载转换为均布荷载的系数取F = 1-2α²+α³。\n\n"
            "重力荷载代表值计算（抗震用）：\n"
            "• 顶层 G6 = 屋面恒载 + 0.5×屋面活载 + 顶层半层墙柱重 = 5028 kN\n"
            "• 标准层 G2~G5 = 楼面恒载 + 0.5×楼面活载 + 上下各半层墙柱重 = 5679 kN\n"
            "• 首层 G1 = 楼面恒载 + 0.5×楼面活载 + 首层半层墙柱重 = 5679 kN\n"
            "• 结构总重力荷载代表值 ΣGi = 33835 kN",
            font_size=12)


# ═══════════════════════════════════════════════════════════
# 幻灯片8：水平地震作用下内力
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "水平地震作用下内力计算", "第4章  底部剪力法 · 反弯点法")
add_page_number(slide, 8)

headers = ["参数", "数值", "参数", "数值"]
rows = [
    ["自振周期 T₁", "0.56 s（顶点位移法）", "地震影响系数 α₁", "0.073"],
    ["底部总剪力 FEK", "1495.51 kN", "顶部附加系数 δn", "0.1148"],
    ["顶部附加地震作用 ΔFn", "171.68 kN", "最大层间位移角", "1/655（<1/550）"],
]
make_table(slide, 1.5, 3.8, [5.5, 4.5, 5.5, 4.5], headers, rows, font_size=12)

add_textbox(slide, 1.5, 7.0, 30, 6.5,
            "计算方法：\n"
            "1. 采用底部剪力法，各楼层水平地震作用 Fi = (Gi·Hi / ΣGj·Hj) × FEK × (1-δn)\n"
            "2. 顶部附加地震作用 ΔFn = δn×FEK = 171.68kN\n"
            "3. D值法分配各柱剪力 Vij = (Dij/ΣDj) × Vi\n"
            "4. 反弯点法确定柱端弯矩（反弯点高度 y = y₀+y₁+y₂+y₃）\n"
            "5. 节点平衡法计算梁端弯矩 Mb = Mc上 + Mc下\n"
            "6. 梁端剪力 Vb = (Ml+Mr)/ln，柱轴力由梁端剪力逐层累加\n\n"
            "最大层间位移角为1/655，小于规范限值1/550，满足要求。",
            font_size=12)


# ═══════════════════════════════════════════════════════════
# 幻灯片9：风荷载作用下内力
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "风荷载作用下内力计算", "第5章  水平风荷载")
add_page_number(slide, 9)

add_textbox(slide, 1.5, 3.5, 30, 0.8,
            "基本风压 W₀ = 0.30 kN/m²，建筑高度18.6m < 30m，不考虑风振系数（βz = 1.0）。",
            font_size=13)

headers = ["楼层", "高度 (m)", "μz", "受风面积 (m²)", "集中力 (kN)", "层间位移角"]
rows = [
    ["6", "19.0", "1.17", "14.49", "8.59", "1/190209"],
    ["5", "16.0", "1.15", "20.70", "12.04", "1/79198"],
    ["4", "13.0", "1.08", "20.70", "11.38", "1/51049"],
    ["3", "10.0", "1.00", "20.70", "10.49", "1/38444"],
    ["2", "7.0", "1.00", "20.70", "10.49", "1/30831"],
    ["1", "4.0", "0.96", "24.15", "11.75", "1/25569"],
]
make_table(slide, 1.5, 5.0, [3, 3.5, 3, 4, 3.5, 4], headers, rows, font_size=11)

add_textbox(slide, 1.5, 11.0, 30, 3.5,
            "风荷载下内力同样采用D值法分配剪力、反弯点法确定柱端弯矩、节点平衡法求梁端弯矩。\n"
            "最大层间位移角为1/25569，远小于规范限值1/550。\n"
            "风荷载产生的内力相比地震作用小很多，不起控制作用。",
            font_size=12)


# ═══════════════════════════════════════════════════════════
# 幻灯片10：竖向荷载下内力计算
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "竖向荷载下内力计算", "第6章  弯矩二次分配法")
add_page_number(slide, 10)

add_textbox(slide, 1.5, 3.5, 30, 0.8,
            "分别计算恒载和活载作用下的框架内力，采用弯矩二次分配法。",
            font_size=13)

steps = [
    "计算节点弯矩分配系数 μik = Sik/ΣSij（与梁柱线刚度成正比）",
    "计算梁端固端弯矩（梯形荷载转等效均布荷载后按两端固定梁公式）",
    "弯矩二次分配：放松节点→分配不平衡弯矩→传递1/2→再次分配并叠加",
    "由最终杆端弯矩反算跨中弯矩：M中 = M简支 - |M左+M右|/2",
    "取隔离体平衡条件计算梁端剪力和柱轴力",
]
add_bullet_list(slide, 1.5, 4.8, 30, steps, font_size=12, spacing=1.3)

add_textbox(slide, 1.5, 13.0, 30, 3.0,
            "边跨梁屋面等效均布荷载16.77kN/m（恒），楼面20.80kN/m（恒）；"
            "中跨梁屋面9.33kN/m（恒），楼面14.64kN/m（恒）。\n"
            "活载下边跨梁等效均布荷载1.44kN/m（屋面）/ 5.74kN/m（楼面）。\n"
            "恒载下顶层边跨梁端弯矩约-66.6kN·m（支座）/ +64.0kN·m（跨中），"
            "各层内力随楼层降低逐渐增大。",
            font_size=12)


# ═══════════════════════════════════════════════════════════
# 幻灯片11：内力组合
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "内力组合", "第7章  无震组合与有震组合")
add_page_number(slide, 11)

add_textbox(slide, 1.5, 3.5, 30, 0.8,
            "依据GB50068-2018，对四种荷载工况进行组合，考虑弯矩调幅和控制截面转换。",
            font_size=13)

headers = ["组合类型", "组合公式", "说明"]
rows = [
    ["无震组合1", "1.3恒 + 1.5活", "可变荷载（活载）控制"],
    ["无震组合2", "1.3恒 + 1.5风", "风荷载控制"],
    ["无震组合3", "1.3恒 + 1.5活 + 1.5×0.6风", "活载为主+风载"],
    ["无震组合4", "1.3恒 + 1.5风 + 1.5×0.7活", "风载为主+活载"],
    ["有震组合", "1.2(恒+0.5活) ± 1.3地震", "地震作用参与"],
]
make_table(slide, 1.5, 4.8, [4.5, 8.5, 5], headers, rows, font_size=11)

add_textbox(slide, 1.5, 10.0, 30, 4.5,
            "处理流程：\n"
            "1. 弯矩调幅：梁端负弯矩乘以0.85（调幅系数），跨中弯矩相应增大\n"
            "2. 控制截面转换：将轴线处内力转换为梁端/跨中控制截面内力\n"
            "   M控 = M轴 - V·b/2（b为柱截面宽度500mm）\n"
            "3. 梁内力组合取各组合的最不利值（|M|max及对应V，Vmax及对应M）\n"
            "4. 柱内力组合按Nmax（轴力最大）和|M|max（弯矩绝对值最大）两组分别处理",
            font_size=12)


# ═══════════════════════════════════════════════════════════
# 幻灯片12：梁截面设计
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "梁截面设计", "第8章  正截面 + 斜截面配筋")
add_page_number(slide, 12)

add_textbox(slide, 1.5, 3.5, 15, 0.8, "设计流程", font_size=16,
            font_color=DARK_BLUE, bold=True)
steps = [
    "强剪弱弯调整：V = ηvb·(Ml+Mr)/ln + VGb",
    "正截面受弯计算：αs = M/(α1·fc·b·h0²)，γs = (1+√1-2αs)/2，As = M/(fy·γs·h0)",
    "最小配筋率验算：As ≥ max(0.20%, 45ft/fy%)·b·h = 312.5mm²",
    "斜截面受剪验算：V ≤ 0.25βc·fc·b·h0",
    "箍筋加密区长度取max(1.5h, 500mm)，间距取100mm",
]
add_bullet_list(slide, 1.5, 4.3, 15, steps, font_size=11, spacing=1.2)

add_textbox(slide, 18, 3.5, 15, 0.8, "典型梁配筋结果", font_size=16,
            font_color=DARK_BLUE, bold=True)
headers = ["楼层", "边跨梁底筋", "边跨支座负筋", "箍筋"]
rows = [
    ["1层", "2C16 (402)", "5C16 (1005)", "C8@100/200"],
    ["3层", "2C16 (402)", "5C16 (1005)", "C8@100/200"],
    ["6层", "2C16 (402)", "5C16 (1005)", "C8@100/200"],
]
make_table(slide, 18, 4.5, [3, 4, 4.5, 4], headers, rows, font_size=11)

add_textbox(slide, 1.5, 12.5, 30, 3.5,
            "说明：表中为mid榀④轴框架梁配筋，括号内为配筋面积mm²。"
            "边跨梁跨中正弯矩按T形截面计算（翼缘宽度bf'=min{ln/3=1633, b+12hf'=1690}=1633mm），"
            "中跨梁跨中以及各支座截面按矩形截面计算。"
            "所有梁截面均满足x ≤ ξb·h0的适筋梁要求，最小配筋率满足规范规定。",
            font_size=11, font_color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════
# 幻灯片13：柱截面设计
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "柱截面设计", "第8章  轴压比验算 + 配筋计算")
add_page_number(slide, 13)

add_textbox(slide, 1.5, 3.5, 15, 0.8, "柱轴压比验算", font_size=14,
            font_color=DARK_BLUE, bold=True)
headers = ["楼层", "边柱轴压比", "中柱轴压比", "限值（三级）"]
rows = [
    ["1层", "0.44", "0.52", "0.85"],
    ["3层", "0.28", "0.34", "0.85"],
    ["6层", "0.06", "0.07", "0.85"],
]
make_table(slide, 1.5, 4.3, [3.5, 3.5, 3.5, 4], headers, rows, font_size=11)

add_textbox(slide, 1.5, 8.5, 15, 4.0,
            "柱配筋设计要点：\n"
            "• 强柱弱梁调整：ΣMc = ηc·ΣMb\n"
            "• 柱端剪力调整：Vc = (Mc上+Mc下)/Hn\n"
            "• 按Nmax和|M|max两组内力分别配筋\n"
            "• 均按对称配筋，单边纵筋4C16(804mm²)\n"
            "• 单边配筋率0.96%，满足最小0.7%（边柱）/0.6%（中柱）\n"
            "• 体积配箍率满足三级抗震构造要求",
            font_size=11)

add_textbox(slide, 18, 3.5, 15, 0.8, "柱配筋结果", font_size=14,
            font_color=DARK_BLUE, bold=True)
headers = ["柱类型", "纵筋", "单边面积率", "箍筋"]
rows = [
    ["边柱", "4C16 (804)", "0.96%", "C8@100/200"],
    ["中柱", "4C16 (804)", "0.96%", "C8@100/200"],
]
make_table(slide, 18, 4.3, [4, 4, 3.5, 4], headers, rows, font_size=11)


# ═══════════════════════════════════════════════════════════
# 幻灯片14：楼板设计
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "楼板设计", "第9章  双向板弹性理论计算")
add_page_number(slide, 14)

add_textbox(slide, 1.5, 3.5, 30, 0.8,
            "取典型A板（边跨双向板）进行设计，板厚120mm，按弹性理论计算。",
            font_size=13)

headers = ["参数", "数值", "参数", "数值"]
rows = [
    ["板短跨 lx", "3.45 m", "板长跨 ly", "5.40 m"],
    ["lx/ly (λ)", "1.57", "计算跨度比", "简支 + 固支"],
    ["恒载设计值", "1.3×4.2 = 5.46 kN/m²", "活载设计值", "1.5×2.0 = 3.00 kN/m²"],
    ["总设计荷载 q", "8.46 kN/m²", "q/2+g 组合", "6.96 kN/m²"],
]
make_table(slide, 1.5, 4.8, [5, 5.5, 5, 5.5], headers, rows, font_size=11)

add_textbox(slide, 1.5, 9.0, 30, 4.5,
            "弯矩计算（考虑泊松比调整 mx(μ) = mx + μ·my）：\n"
            "• 跨中短边弯矩：4.50 kN·m/m\n"
            "• 跨中长边弯矩：2.22 kN·m/m\n"
            "• 支座短边负弯矩：-8.28 kN·m/m\n"
            "• 支座长边负弯矩：-5.51 kN·m/m\n\n"
            "配筋结果（楼板按C8@200双层双向配筋，实配面积251mm²/m）：\n"
            "均满足最小配筋率0.20%（240mm²/m）要求。",
            font_size=12)


# ═══════════════════════════════════════════════════════════
# 幻灯片15：楼梯设计
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "楼梯设计", "第10章  板式楼梯")
add_page_number(slide, 15)

add_textbox(slide, 1.5, 3.5, 15, 0.8, "设计参数", font_size=14,
            font_color=DARK_BLUE, bold=True)
params = [
    "楼梯间开间 3.3m，梯段宽 1.5m",
    "踏步尺寸 167×280mm（高×宽）",
    "梯段板厚度 100mm",
    "休息平台宽 1.60m，平台板厚 100mm（取80mm）",
    "梯梁截面 200×350mm，跨度 3m",
]
add_bullet_list(slide, 1.5, 4.3, 15, params, font_size=12, spacing=1.1)

add_textbox(slide, 18, 3.5, 15, 0.8, "计算结果", font_size=14,
            font_color=DARK_BLUE, bold=True)
results = [
    "踏步板恒载：7.33kN/m²，设计值14.52kN/m²",
    "踏步板配筋：跨中C8@150(335mm²)，支座C8@200(251mm²)",
    "平台板恒载：3.92kN/m²，设计值10.35kN/m²",
    "平台板配筋：C8@200 (251mm²)",
    "梯梁设计荷载27.81kN/m",
    "梯梁配筋：底2C16，箍筋C8@100/200",
]
add_bullet_list(slide, 18, 4.3, 15, results, font_size=12, spacing=1.1)

add_textbox(slide, 1.5, 12.5, 30, 3.0,
            "踏步板按斜向简支板计算：斜向跨度=2.696m，M=0.1ql²=7.29kN·m。\n"
            "梯梁两端简支于框架梁上，按简支梁计算跨中弯矩后配筋。\n"
            "楼梯各参数取自梯间几何尺寸，不受框架横向跨度变化影响。",
            font_size=11, font_color=DARK_GRAY)


# ═══════════════════════════════════════════════════════════
# 幻灯片16：基础设计
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "基础设计", "第11章  柱下独立基础")
add_page_number(slide, 16)

add_textbox(slide, 1.5, 3.5, 30, 0.8,
            "地基持力层粉质黏土，fak=180kPa，修正后fa=199.24kPa。边柱独立基础，中柱联合基础。",
            font_size=13)

headers = ["项目", "边柱独立基础", "中柱联合基础"]
rows = [
    ["基底尺寸", "B×L = 2.5m × 3.6m (9.0m²)", "B×L = 3.6m × 5.6m (20.2m²)"],
    ["基础高度", "h = 750mm", "h = 750mm"],
    ["轴心基底压力", "178.4 kPa < fa", "195.0 kPa < fa"],
    ["偏心基底压力", "181.1 kPa < 1.2fa", "满足要求"],
    ["抗冲切验算", "Fl=357.9kN < 910.9kN", "Fl=227.4kN < 910.9kN"],
    ["底板配筋（沿长边）", "C12@100 (1131mm²/m)", "C18@200 (1272mm²/m)"],
    ["底板配筋（沿短边）", "C12@100 (1131mm²/m)", "C12@100 (1131mm²/m)"],
]
make_table(slide, 1.5, 4.8, [4.5, 6.5, 7], headers, rows, font_size=10)

add_textbox(slide, 1.5, 12.0, 30, 3.5,
            "设计流程：\n"
            "1. 取底层柱底标准组合内力（边柱Nk=1269.9kN，中柱Nk=1612.5kN）\n"
            "2. 按轴心受压初步确定基底面积 A ≥ Nk/(fa-γm·d)\n"
            "3. 验算基底压力（含偏心荷载）pk ≤ fa，pkmax ≤ 1.2fa\n"
            "4. 抗冲切验算 Fl ≤ 0.7·βhp·ft·am·h0\n"
            "5. 底板按悬臂板计算弯矩，按最小配筋率0.15%及计算值取大配筋",
            font_size=12)


# ═══════════════════════════════════════════════════════════
# 幻灯片17：结论
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
shape = slide.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Cm(0), Cm(0), Cm(SLD_W), Cm(SLD_H))
shape.fill.solid()
shape.fill.fore_color.rgb = DARK_BLUE
shape.line.fill.background()

add_textbox(slide, 2, 3.0, 30, 2.0, "设计总结",
            font_size=36, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)

conclusions = [
    "本设计完成了6层钢筋混凝土框架结构的全套结构计算",
    "内容涵盖荷载统计、内力分析、截面配筋到基础设计的全流程",
    "抗震设计采用底部剪力法，竖向内力采用弯矩二次分配法",
    "楼盖采用双向板弹性理论，楼梯为板式楼梯",
    "所有构件均满足强度、刚度及抗震构造要求",
    "计算书共计11章，包含80个计算表格",
]
for i, c in enumerate(conclusions):
    add_textbox(slide, 5, 7.0 + i * 1.3, 25, 1.0, f"  {c}",
                font_size=16, font_color=RGBColor(0xDD, 0xE8, 0xF0))

add_textbox(slide, 2, 15.5, 30, 1.5, "恳请各位老师批评指正！",
            font_size=28, font_color=WHITE, bold=True,
            alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════
# 幻灯片18：参考文献
# ═══════════════════════════════════════════════════════════
slide = add_blank_slide()
add_title_bar(slide, "参考文献")
add_page_number(slide, 18)

refs = [
    "[1] GB50068-2018, 建筑结构可靠性设计统一标准[S]. 北京: 中国建筑工业出版社, 2018.",
    "[2] GB50009-2012, 建筑结构荷载规范[S]. 北京: 中国建筑工业出版社, 2012.",
    "[3] GB/T50011-2010, 建筑抗震设计规范(2024版)[S]. 北京: 中国建筑工业出版社, 2024.",
    "[4] GB/T50010-2010, 混凝土结构设计规范(2024版)[S]. 北京: 中国建筑工业出版社, 2024.",
    "[5] GB50007-2011, 建筑地基基础设计规范[S]. 北京: 中国建筑工业出版社, 2011.",
    "[6] GB55001-2021, 工程结构通用规范[S]. 北京: 中国建筑工业出版社, 2021.",
    "[7] GB55008-2021, 混凝土结构通用规范[S]. 北京: 中国建筑工业出版社, 2021.",
    "[8] 梁兴文. 混凝土结构设计原理(第五版)[M]. 北京: 中国建筑工业出版社, 2022.",
    "[9] 梁兴文. 混凝土结构设计(第五版)[M]. 北京: 中国建筑工业出版社, 2022.",
    "[10] 梁兴文, 史庆轩. 土木工程专业毕业设计指导-房屋建筑工程卷[M]. 北京: 中国建筑工业出版社, 2014.",
    "[11] 王社良. 抗震结构设计(第5版)[M]. 北京: 中国建筑工业出版社, 2022.",
]
add_bullet_list(slide, 1.5, 3.8, 30, refs, font_size=12, spacing=1.2)


# ═══════════════════════════════════════════════════════════
# 保存
# ═══════════════════════════════════════════════════════════
output_path = r"C:\Users\邓杰鹏\Desktop\毕设\答辩PPT_邓杰鹏.pptx"
prs.save(output_path)
print(f"PPT已保存至: {output_path}")
print(f"共 {len(prs.slides)} 页")
