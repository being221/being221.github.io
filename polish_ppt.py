"""
PPT润色：补空白 + 删风险点 + 加答辩引导
"""
import sys,io;sys.stdout=io.TextIOWrapper(sys.stdout.buffer,encoding='utf-8')
from pptx import Presentation
from pptx.util import Pt,Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

INPUT=r'C:/Users/邓杰鹏/Desktop/毕设答辩-邓杰鹏-终版.pptx'
prs=Presentation(INPUT)
SW=prs.slide_width;SH=prs.slide_height

def fill(shape,lines):
    tf=shape.text_frame
    for i,(t,sz,b) in enumerate(lines):
        p=tf.paragraphs[i] if i<len(tf.paragraphs) else tf.add_paragraph()
        p.clear();r=p.add_run();r.text=t;r.font.size=Pt(sz)
        if b is not None:r.font.bold=b
    for j in range(len(lines),len(tf.paragraphs)):
        tf.paragraphs[j].clear()

def add_page_number(slide,num):
    """右下角页码"""
    txBox=slide.shapes.add_textbox(Emu(11000000),Emu(6600000),Emu(1000000),Emu(400000))
    tf=txBox.text_frame
    p=tf.paragraphs[0];p.alignment=PP_ALIGN.RIGHT
    r=p.add_run();r.text=str(num);r.font.size=Pt(10)
    from pptx.dml.color import RGBColor
    r.font.color.rgb=RGBColor(0x99,0x99,0x99)

# ═══════════════════════════════════════════════════════
# P03: 章节封面 — 加一个简短的本章要点
# ═══════════════════════════════════════════════════════
print('P03: 补充章节指引')
slide=prs.slides[2]
# 在底部加一行指引文字
txBox=slide.shapes.add_textbox(Emu(3400000),Emu(5000000),Emu(9000000),Emu(1600000))
tf=txBox.text_frame;tf.word_wrap=True
p=tf.paragraphs[0]
r=p.add_run();r.text='本章内容';r.font.size=Pt(14);r.font.bold=True
r.font.color.rgb=RGBColor(0x1A,0x23,0x7E)
p=tf.add_paragraph();p.space_before=Pt(4)
r=p.add_run()
r.text='项目背景 · 结构方案 · 设计参数 · 材料选用'
r.font.size=Pt(12)

# ═══════════════════════════════════════════════════════
# P04: 加答辩引导 — "设计依据"
# ═══════════════════════════════════════════════════════
print('P04: 加设计依据引导')
slide=prs.slides[3]
for s in slide.shapes:
    if s.has_text_frame and '项目概况' in s.text_frame.text:
        tf=s.text_frame
        # 在末尾追加一行引用
        p=tf.add_paragraph()
        p.space_before=Pt(8)
        r=p.add_run()
        r.text='规范依据：GB50011 抗规 · GB50009 荷规 · GB50010 混规'
        r.font.size=Pt(10);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P05: 竖向荷载 — 补"答辩要点"提示（对用户不可见，对老师安全）
# ═══════════════════════════════════════════════════════
print('P05: 补设计流程图解')
slide=prs.slides[4]
# 已够充实，加上弯矩分配法简短说明
for s in slide.shapes:
    if s.has_text_frame and '第四步' in s.text_frame.text:
        tf=s.text_frame
        p=tf.add_paragraph()
        p.space_before=Pt(6)
        r=p.add_run()
        r.text='活载计算方法相同，仅荷载值不同。详见计算书第6章。'
        r.font.size=Pt(11);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P06: 地震 — 已充实
# ═══════════════════════════════════════════════════════
print('P06: 已充实')

# ═══════════════════════════════════════════════════════
# P07: 风荷载 — 表格下方加简注
# ═══════════════════════════════════════════════════════
print('P07: 加表格注释')
slide=prs.slides[6]
for s in slide.shapes:
    if s.has_text_frame and '风荷载内力计算' in s.text_frame.text:
        tf=s.text_frame
        p=tf.add_paragraph()
        r=p.add_run()
        r.text='各层风荷载集中力作用于楼层标高处，按D值法分配至各柱'
        r.font.size=Pt(10);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P08: D值法 — 补答辩引导
# ═══════════════════════════════════════════════════════
print('P08: 补答辩引导')
slide=prs.slides[7]
for s in slide.shapes:
    if s.has_text_frame and '方法要点' in s.text_frame.text:
        tf=s.text_frame
        p=tf.add_paragraph()
        p.space_before=Pt(6)
        r=p.add_run()
        r.text='D值法本质：修正反弯点法，考虑了梁柱线刚度比和楼层位置的影响'
        r.font.size=Pt(10);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P09: 荷载组合 — 补背景说明
# ═══════════════════════════════════════════════════════
print('P09: 补组合背景')
slide=prs.slides[8]
for s in slide.shapes:
    if s.has_text_frame and '组合目的' in s.text_frame.text:
        tf=s.text_frame
        p=tf.add_paragraph()
        p.space_before=Pt(6)
        r=p.add_run()
        r.text='依据GB50009第3.2节，考虑荷载同时出现的概率，取不同分项系数组合'
        r.font.size=Pt(10);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P11: 梁截面设计 — 57%利用率，补答辩引导
# ═══════════════════════════════════════════════════════
print('P11: 填充梁设计内容')
slide=prs.slides[10]
for s in slide.shapes:
    if s.has_text_frame and '构造配箍' in s.text_frame.text:
        tf=s.text_frame
        # 加抗震构造说明
        p=tf.add_paragraph()
        p.space_before=Pt(6)
        r=p.add_run()
        r.text='抗震要求：梁端箍筋加密区长度≥1.5h=750mm，间距≤min(h/4,8d,100)'
        r.font.size=Pt(10);r.font.italic=True
        # 加T形截面补充
        p=tf.add_paragraph()
        p.space_before=Pt(2)
        r=p.add_run()
        r.text='T形截面翼缘宽度取min(l₀/3, b+sn, b+12hf\')=1600mm'
        r.font.size=Pt(10);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P12: 柱截面设计 — 补轴压比具体值
# ═══════════════════════════════════════════════════════
print('P12: 补轴压比数据')
slide=prs.slides[11]
for s in slide.shapes:
    if s.has_text_frame and '轴压比' in s.text_frame.text:
        tf=s.text_frame
        p=tf.add_paragraph()
        p.space_before=Pt(4)
        r=p.add_run()
        r.text='以一层边柱为例：Nmax=1727kN, 轴压比=1727×10³/(14.3×500×500)=0.48 < 0.85'
        r.font.size=Pt(10);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P13: 柱公式 — 去"小偏心"残留文字（模板遗留）
# ═══════════════════════════════════════════════════════
print('P13: 清理残留')
slide=prs.slides[12]
for s in list(slide.shapes):
    if s.has_text_frame:
        if '柱子钢筋的构造' in s.text_frame.text:
            s._element.getparent().remove(s._element)

# ═══════════════════════════════════════════════════════
# P15: 基础 — 补答辩引导
# ═══════════════════════════════════════════════════════
print('P15: 补答辩引导')
slide=prs.slides[14]
for s in slide.shapes:
    if s.has_text_frame and '基础选型' in s.text_frame.text:
        tf=s.text_frame
        p=tf.add_paragraph()
        p.space_before=Pt(4)
        r=p.add_run()
        r.text='中柱采用联合基础原因：柱距较近(2.4m)，单基础重叠，合并更经济'
        r.font.size=Pt(10);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P16: 设计成果 — 已充实，加规范汇总
# ═══════════════════════════════════════════════════════
print('P16: 加规范')
slide=prs.slides[15]
for s in slide.shapes:
    if s.has_text_frame and '结构方案' in s.text_frame.text:
        tf=s.text_frame
        p=tf.add_paragraph()
        p.space_before=Pt(6)
        r=p.add_run()
        r.text='设计遵循规范：GB50011 · GB50009 · GB50010 · GB50007 · GB50068'
        r.font.size=Pt(10);r.font.italic=True

# ═══════════════════════════════════════════════════════
# P18: 致谢 — 补答辩收尾引导
# ═══════════════════════════════════════════════════════
print('P18: 补收尾')
slide=prs.slides[17]
for s in slide.shapes:
    if s.has_text_frame and '感谢' in s.text_frame.text:
        tf=s.text_frame
        p=tf.add_paragraph()
        p.space_before=Pt(12)
        r=p.add_run()
        r.text='请各位老师提问，谢谢！'
        r.font.size=Pt(22);r.font.bold=True

# ═══════════════════════════════════════════════════════
# 全局：所有页加页码
# ═══════════════════════════════════════════════════════
print('添加页码...')
for si,slide in enumerate(prs.slides):
    num=str(si+1)
    txBox=slide.shapes.add_textbox(Emu(11300000),Emu(6600000),Emu(700000),Emu(300000))
    tf=txBox.text_frame;tf.word_wrap=False
    p=tf.paragraphs[0];p.alignment=PP_ALIGN.RIGHT
    r=p.add_run();r.text=num;r.font.size=Pt(9)
    r.font.color.rgb=RGBColor(0x99,0x99,0x99)

# ═══════════════════════════════════════════════════════
# 最终验证
# ═══════════════════════════════════════════════════════
print('\n===== 最终验证 =====')
issues=0
for si,slide in enumerate(prs.slides):
    n=si+1
    si_issues=[]
    for s in slide.shapes:
        if s.left/360000>35 or s.top/360000>35:continue
        r=s.left+s.width;b=s.top+s.height
        if s.left>Emu(-100000) and s.top>Emu(-100000):
            if r>SW+Emu(50000):si_issues.append('overflow')
            if b>SH+Emu(80000):si_issues.append('overflow')
    main=''
    for s in slide.shapes:
        if s.has_text_frame:
            t=s.text_frame.text.replace('\n',' ').strip()[:70]
            if t:main=t;break
    ok=' OK' if not si_issues else ' FIX'
    if si_issues:issues+=1
    print(f'  P{n:02d}: {main}{ok}')
print(f'\nIssues: {issues}')

prs.save(INPUT)
print('Saved!')
